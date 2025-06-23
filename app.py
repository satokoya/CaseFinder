import os
import sqlite3
from flask import Flask, render_template, request, redirect, url_for, flash, send_file
from werkzeug.utils import secure_filename
from pptx import Presentation
import PyPDF2
from openpyxl import load_workbook
from io import BytesIO
from flask_httpauth import HTTPBasicAuth
from werkzeug.security import generate_password_hash, check_password_hash
from datetime import datetime

# アプリケーションの初期化
app = Flask(__name__)
app.secret_key = 'dev'  # 本番環境では強力なシークレットキーに変更

# 認証の設定
auth = HTTPBasicAuth()

# ユーザー名とパスワードの設定（本番環境ではデータベースなどに保存）
users = {
    "admin": generate_password_hash("sangyokikaku")  # パスワードを変更してください
}

@auth.verify_password
def verify_password(username, password):
    if username in users and check_password_hash(users.get(username), password):
        return username

@auth.error_handler
def auth_error(status):
    return "認証に失敗しました。正しいユーザー名とパスワードを入力してください。", status

# デバッグ用
print("アプリケーションが起動しました")
print(f"認証が有効化されています: {auth is not None}")

# アップロードされたファイルを保存するディレクトリ
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# 許可するファイル拡張子
ALLOWED_EXTENSIONS = {'pptx', 'pdf', 'xlsx', 'xls'}

# SQLite 設定
DB_PATH = 'casefinder.db'

def init_db():
    conn = sqlite3.connect(DB_PATH)
    c = conn.cursor()
    
    # Check if cases table exists
    c.execute("SELECT name FROM sqlite_master WHERE type='table' AND name='cases'")
    table_exists = c.fetchone()
    
    if not table_exists:
        # Create new tables with created_at column
        c.execute(
            """
            CREATE TABLE cases (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                filename TEXT NOT NULL,
                stored_path TEXT NOT NULL,
                text_content TEXT NOT NULL,
                customer_name TEXT,
                system_name TEXT,
                created_at DATETIME DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        conn.commit()
    conn.close()

# Initialize database with the latest schema
init_db()

def get_db_connection():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    """Extract text content from PDF file."""
    text = []
    try:
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text.append(page.extract_text() or '')
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""

def extract_text_from_excel(file_path):
    """Extract text content from Excel file."""
    text = []
    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:  # Only add non-empty rows
                    text.append(' | '.join(row_text))
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return ""

def extract_text_from_pptx(file_path):
    """Extract text content from PowerPoint file."""
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""

@app.route('/')
@auth.login_required
def index():
    conn = get_db_connection()
    cases = conn.execute('SELECT * FROM cases ORDER BY created_at DESC').fetchall()
    conn.close()
    return render_template('index.html', cases=cases)

@app.route('/upload', methods=['POST'])
@auth.login_required
def upload():
    if 'file' not in request.files:
        flash('ファイルが選択されていません')
        return redirect(request.url)
    
    file = request.files['file']
    if file.filename == '':
        flash('ファイルが選択されていません')
        return redirect(request.url)
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # ファイルが既に存在する場合はリネーム
        counter = 1
        name, ext = os.path.splitext(filename)
        while os.path.exists(filepath):
            filename = f"{name}_{counter}{ext}"
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            counter += 1
        
        file.save(filepath)
        
        # ファイルの種類に応じてテキストを抽出
        if filename.lower().endswith('.pptx'):
            text_content = extract_text_from_pptx(filepath)
        elif filename.lower().endswith('.pdf'):
            text_content = extract_text_from_pdf(filepath)
        elif filename.lower().endswith(('.xlsx', '.xls')):
            text_content = extract_text_from_excel(filepath)
        else:
            text_content = ""
        
        # データベースに保存
        conn = get_db_connection()
        conn.execute(
            'INSERT INTO cases (filename, stored_path, text_content) VALUES (?, ?, ?)',
            (filename, filepath, text_content)
        )
        conn.commit()
        conn.close()
        
        flash('ファイルが正常にアップロードされました')
        return redirect(url_for('index'))
    
    flash('許可されていないファイル形式です')
    return redirect(request.url)

@app.route('/case/<int:case_id>')
@auth.login_required
def case_detail(case_id):
    conn = get_db_connection()
    case = conn.execute('SELECT * FROM cases WHERE id = ?', (case_id,)).fetchone()
    conn.close()
    
    if case is None:
        flash('指定されたケースは見つかりませんでした')
        return redirect(url_for('index'))
    
    return render_template('case_detail.html', case=case)

@app.route('/update_metadata/<int:case_id>', methods=['POST'])
@auth.login_required
def update_metadata(case_id):
    customer_name = request.form.get('customer_name', '')
    system_name = request.form.get('system_name', '')
    
    conn = get_db_connection()
    conn.execute(
        'UPDATE cases SET customer_name = ?, system_name = ? WHERE id = ?',
        (customer_name, system_name, case_id)
    )
    conn.commit()
    conn.close()
    
    flash('メタデータが更新されました')
    return redirect(url_for('case_detail', case_id=case_id))

@app.route('/download_pdf/<int:case_id>')
@auth.login_required
def download_pdf(case_id):
    conn = get_db_connection()
    case = conn.execute('SELECT * FROM cases WHERE id = ?', (case_id,)).fetchone()
    conn.close()
    
    if case is None:
        flash('指定されたケースは見つかりませんでした')
        return redirect(url_for('index'))
    
    # PDFの作成
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    
    # 日本語フォントの登録
    pdfmetrics.registerFont(UnicodeCIDFont('HeiseiKakuGo-W5'))
    
    buffer = BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    # タイトル
    title = f"{case['customer_name']}_{case['system_name']}" if case['customer_name'] and case['system_name'] else f"事例要約: {case['filename']}"
    p.setFont('HeiseiKakuGo-W5', 16)
    p.drawString(50, height - 50, title)
    
    # ファイル情報
    p.setFont('HeiseiKakuGo-W5', 12)
    p.drawString(50, height - 80, f"ファイル名: {case['filename']}")
    p.drawString(50, height - 100, f"アップロード日時: {case['created_at']}")
    
    # テキストコンテンツ
    y_position = height - 140
    p.setFont('HeiseiKakuGo-W5', 10)
    for line in case['text_content'].split('\n'):
        if y_position < 50:  # 新しいページを追加
            p.showPage()
            y_position = height - 50
            p.setFont('HeiseiKakuGo-W5', 10)
        p.drawString(50, y_position, line)
        y_position -= 15
    
    p.save()
    
    buffer.seek(0)
    return send_file(
        buffer,
        as_attachment=True,
        download_name=f"{os.path.splitext(case['filename'])[0]}_summary.pdf",
        mimetype='application/pdf'
    )

@app.route('/delete/<int:case_id>', methods=['POST'])
@auth.login_required
def delete_case(case_id):
    conn = get_db_connection()
    case = conn.execute('SELECT * FROM cases WHERE id = ?', (case_id,)).fetchone()
    
    if case is None:
        flash('指定されたケースは見つかりませんでした')
        return redirect(url_for('index'))
    
    # ファイルを削除
    try:
        if os.path.exists(case['stored_path']):
            os.remove(case['stored_path'])
    except Exception as e:
        print(f"Error deleting file: {e}")
    
    # データベースから削除
    conn.execute('DELETE FROM cases WHERE id = ?', (case_id,))
    conn.commit()
    conn.close()
    
    flash('ケースが削除されました')
    return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(debug=True)